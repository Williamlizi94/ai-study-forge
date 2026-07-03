from __future__ import annotations

import re
from io import BytesIO
from pathlib import Path

from docx import Document
from pptx import Presentation
from pypdf import PdfReader


SUPPORTED_EXTENSIONS = {".txt", ".pdf", ".docx", ".doc", ".pptx", ".ppt"}


class DocumentParseError(ValueError):
    pass


def extract_text_from_document(filename: str, content: bytes) -> str:
    extension = Path(filename).suffix.lower()
    if extension not in SUPPORTED_EXTENSIONS:
        supported = ", ".join(sorted(SUPPORTED_EXTENSIONS))
        raise DocumentParseError(f"Unsupported file type. Supported types: {supported}.")

    if not content:
        raise DocumentParseError("Uploaded file is empty.")

    if extension == ".txt":
        text = _extract_txt(content)
    elif extension == ".pdf":
        text = _extract_pdf(content)
    elif extension == ".docx":
        text = _extract_docx(content)
    elif extension == ".pptx":
        text = _extract_pptx(content)
    else:
        text = _extract_legacy_binary_best_effort(content)

    text = _clean_text(text)
    if len(text) < 50:
        raise DocumentParseError(
            "Could not extract enough readable text from this file. "
            "For scanned PDFs, use OCR first. For legacy .doc/.ppt files, try saving as .docx/.pptx."
        )
    return text


def _extract_txt(content: bytes) -> str:
    for encoding in ("utf-8-sig", "utf-8", "gb18030", "cp1252"):
        try:
            return content.decode(encoding)
        except UnicodeDecodeError:
            continue
    return content.decode("utf-8", errors="ignore")


def _extract_pdf(content: bytes) -> str:
    try:
        reader = PdfReader(BytesIO(content))
    except Exception as exc:
        raise DocumentParseError("Could not read this PDF file.") from exc

    pages: list[str] = []
    for index, page in enumerate(reader.pages, start=1):
        try:
            page_text = page.extract_text() or ""
        except Exception:
            page_text = ""
        if page_text.strip():
            pages.append(f"Page {index}\n{page_text}")

    if not pages:
        raise DocumentParseError(
            "No selectable text was found in this PDF. Scanned PDFs need OCR before upload."
        )
    return "\n\n".join(pages)


def _extract_docx(content: bytes) -> str:
    try:
        document = Document(BytesIO(content))
    except Exception as exc:
        raise DocumentParseError("Could not read this DOCX file.") from exc

    parts: list[str] = []
    parts.extend(paragraph.text for paragraph in document.paragraphs if paragraph.text.strip())

    for table in document.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if row_text:
                parts.append(row_text)

    return "\n".join(parts)


def _extract_pptx(content: bytes) -> str:
    try:
        presentation = Presentation(BytesIO(content))
    except Exception as exc:
        raise DocumentParseError("Could not read this PPTX file.") from exc

    slides: list[str] = []
    for index, slide in enumerate(presentation.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    parts.append(text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                    if row_text:
                        parts.append(row_text)
        if parts:
            slides.append(f"Slide {index}\n" + "\n".join(parts))

    if not slides:
        raise DocumentParseError("No readable text was found in this PPTX file.")
    return "\n\n".join(slides)


def _extract_legacy_binary_best_effort(content: bytes) -> str:
    # Legacy .doc/.ppt are binary formats. This extracts readable text sequences
    # without requiring LibreOffice or a dedicated conversion service.
    unicode_chunks = []
    for match in re.finditer(rb"(?:[\x20-\x7e]\x00|[\x09\x0a\x0d]\x00){8,}", content):
        chunk = match.group(0).decode("utf-16le", errors="ignore")
        unicode_chunks.append(chunk)

    ascii_chunks = []
    for match in re.finditer(rb"[\x09\x0a\x0d\x20-\x7e]{24,}", content):
        chunk = match.group(0).decode("cp1252", errors="ignore")
        ascii_chunks.append(chunk)

    return "\n".join(unicode_chunks + ascii_chunks)


def _clean_text(text: str) -> str:
    text = text.replace("\x00", " ")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    lines = [line.strip() for line in text.splitlines()]
    return "\n".join(line for line in lines if line)
